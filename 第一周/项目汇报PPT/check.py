from pptx import Presentation
p = Presentation(r'D:\obsidian1\第一周\项目汇报PPT\第一周项目汇报.pptx')
print('slides=', len(p.slides._sldIdLst))
for i, sl in enumerate(p.slides, 1):
    sh = len(sl.shapes)
    has_chart = any(s.has_chart for s in sl.shapes if hasattr(s, 'has_chart'))
    print(f'  p{i:02d} shapes={sh} chart={has_chart}')
