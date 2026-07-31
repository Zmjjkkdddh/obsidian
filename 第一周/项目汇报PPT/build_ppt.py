# -*- coding: utf-8 -*-
"""第一周项目汇报 PPT 生成脚本（python-pptx 1.0.2）
主题：AI 工具上手与 WorkBuddy 自动化实践
受众：管理层 | 风格：商务浅色 | 页数：16
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---------- 设计系统 ----------
NAVY   = "1F3A5F"   # 主色 深蓝
BLUE   = "2E6FB7"   # 强调蓝
ACCENT = "F2A93B"   # 强调琥珀（数字）
GRAY   = "5A6472"   # 辅助灰
LIGHT  = "DCE3EC"   # 浅分隔线
BG     = "FFFFFF"   # 页面背景
SOFT   = "F4F6F9"   # 浅底
DARK   = "232A33"   # 正文深
FONT   = "Microsoft YaHei"

TOTAL = 16
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor.from_string(BG)
    return s

def rect(s, x, y, w, h, fill=None, line=None, lw=1, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def oval(s, x, y, w, h, fill=None, line=None, lw=1):
    return rect(s, x, y, w, h, fill, line, lw, MSO_SHAPE.OVAL)

def text(s, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paragraphs: list of dicts {runs:[(t,sz,col,bold)], align, space_after, line_spacing, level}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', align)
        p.space_after = Pt(para.get('space_after', 4))
        p.space_before = Pt(para.get('space_before', 0))
        if 'line_spacing' in para: p.line_spacing = para['line_spacing']
        for (t, sz, col, bd) in para.get('runs', [('', 14, DARK, False)]):
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = bd
            r.font.color.rgb = RGBColor.from_string(col)
            r.font.name = FONT
    return tb

def title_bar(s, title, kicker=None):
    rect(s, 0.6, 0.55, 0.13, 0.62, fill=ACCENT)
    text(s, 0.85, 0.46, 11.6, 0.8,
         [{'runs': [(title, 26, NAVY, True)], 'anchor': 'mid'}], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, 0.87, 1.18, 11.5, 0.3,
             [{'runs': [(kicker, 12, BLUE, True)]}])
    rect(s, 0.6, 1.40, 12.13, 0.018, fill=LIGHT)

def footer(s, n):
    text(s, 0.6, 7.04, 6.0, 0.3,
         [{'runs': [('第一周项目汇报 · AI 工具上手与 WorkBuddy 实践', 9, GRAY, False)]}])
    text(s, 9.3, 7.04, 3.43, 0.3,
         [{'runs': [(f'{n:02d} / {TOTAL}', 9, GRAY, False)]}], align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, big, label, sub=None, accent=BLUE, big_sz=38):
    rect(s, x, y, w, h, fill=BG, line=LIGHT, lw=1)
    rect(s, x, y, w, 0.10, fill=accent)
    text(s, x, y + 0.30, w, 0.95,
         [{'runs': [(big, big_sz, accent, True)], 'align': PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.12, y + 1.22, w - 0.24, 0.45,
         [{'runs': [(label, 13, NAVY, True)], 'align': PP_ALIGN.CENTER}])
    if sub:
        text(s, x + 0.12, y + 1.66, w - 0.24, 0.5,
             [{'runs': [(sub, 10, GRAY, False)], 'align': PP_ALIGN.CENTER, 'line_spacing': 1.0}])

# =====================================================================
# Slide 1 — 封面
# =====================================================================
s = slide()
rect(s, 0, 0, 0.28, H, fill=NAVY)            # 左侧书脊
oval(s, 10.4, -1.2, 3.6, 3.6, fill=SOFT)     # 装饰圆
oval(s, 11.6, 1.4, 1.5, 1.5, fill="EAF1F8")
text(s, 1.0, 2.35, 11.0, 0.5,
     [{'runs': [('PROJECT REVIEW · 2026', 13, BLUE, True)]}])
text(s, 1.0, 2.85, 11.4, 1.3,
     [{'runs': [('第一周项目汇报', 46, NAVY, True)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.02, 4.15, 11.4, 0.7,
     [{'runs': [('AI 工具上手与 WorkBuddy 自动化实践', 22, GRAY, False)]}])
rect(s, 1.02, 5.0, 3.2, 0.06, fill=ACCENT)
text(s, 1.02, 6.1, 11.4, 0.9,
     [{'runs': [
        ('汇报对象：管理层    ', 13, DARK, False),
        ('|    汇报周期：2026/07/28 – 07/30    ', 13, DARK, False),
        ('|    汇报日期：2026/07/31', 13, DARK, False)]}])

# =====================================================================
# Slide 2 — 核心结论
# =====================================================================
s = slide()
title_bar(s, '核心结论', 'EXECUTIVE SUMMARY')
text(s, 0.85, 1.62, 11.6, 0.55,
     [{'runs': [('一周内完成 AI 工具体系从 0 到 1 搭建，跑通自动化复盘闭环，验证可复制的提效路径。', 14, BLUE, True)]}])
xs = [0.6, 3.65, 6.7, 9.75]; cw = 2.85
data = [
    ('01', '知识库成型', '累计 25 页结构化笔记，覆盖 10 大 AI 概念与 5 类工具。', BLUE),
    ('02', '自动化跑通', 'WorkBuddy 工作流每月自动生成知识库可视化报告，复盘自动化。', ACCENT),
    ('03', '工具实操', '完成 Claude Code / Codex 实测与对比，产出可运行 Demo。', BLUE),
    ('04', '方法论沉淀', '引入 Karpathy LLM Wiki 知识管理法，具备团队推广基础。', ACCENT),
]
for x, (no, t, d, ac) in zip(xs, data):
    rect(s, x, 2.35, cw, 2.5, fill=BG, line=LIGHT, lw=1)
    rect(s, x, 2.35, cw, 0.10, fill=ac)
    text(s, x + 0.18, 2.55, cw - 0.3, 0.5, [{'runs': [(no, 22, ac, True)]}])
    text(s, x + 0.18, 3.08, cw - 0.3, 0.45, [{'runs': [(t, 16, NAVY, True)]}])
    text(s, x + 0.18, 3.6, cw - 0.36, 1.1,
         [{'runs': [(d, 12, GRAY, False)], 'line_spacing': 1.1}])
# 底部数字条
rect(s, 0.6, 5.2, 12.13, 1.05, fill=NAVY)
nums = [('25', '累计笔记页数'), ('3', '集中实践天数'), ('1', '自动化工作流'), ('1', '可复用方法论')]
for i, (b, l) in enumerate(nums):
    nx = 0.6 + i * 3.03
    text(s, nx, 5.32, 3.03, 0.6, [{'runs': [(b, 30, ACCENT, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, nx, 5.9, 3.03, 0.3, [{'runs': [(l, 12, 'FFFFFF', False)], 'align': PP_ALIGN.CENTER}])
    if i < 3:
        rect(s, nx + 3.03, 5.45, 0.015, 0.55, fill="3A567D")
footer(s, 2)

# =====================================================================
# Slide 3 — 汇报框架 / 目录
# =====================================================================
s = slide()
title_bar(s, '汇报框架', 'AGENDA')
ag = [
    ('01', '背景与目标', '为什么做、做到什么程度', BLUE),
    ('02', '周进展与成果', '三天时间线与三大核心产出', ACCENT),
    ('03', '经验与方法论', 'LLM Wiki 知识管理法 + 问题解决', BLUE),
    ('04', '下一步与所需支持', '推广计划、风险与资源请求', ACCENT),
]
for i, (no, t, d, ac) in enumerate(ag):
    y = 1.85 + i * 1.18
    rect(s, 0.6, y, 12.13, 1.0, fill=SOFT)
    rect(s, 0.6, y, 0.12, 1.0, fill=ac)
    text(s, 0.95, y, 1.2, 1.0, [{'runs': [(no, 30, ac, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.25, y + 0.16, 7.0, 0.45, [{'runs': [(t, 18, NAVY, True)]}])
    text(s, 2.25, y + 0.6, 9.0, 0.35, [{'runs': [(d, 12, GRAY, False)]}])
footer(s, 3)

# =====================================================================
# Slide 4 — 背景与目标
# =====================================================================
s = slide()
title_bar(s, '项目背景与目标', 'BACKGROUND & OBJECTIVES')
text(s, 0.85, 1.65, 11.6, 1.4,
     [{'runs': [
        ('背景：', 14, NAVY, True),
        ('AI 工具与编码智能体（Agent）迭代迅速，团队亟需快速建立系统认知与实操能力，把"会用"升级为"能沉淀、可复制"。', 14, DARK, False)],
       'space_after': 8, 'line_spacing': 1.2},
      {'runs': [
        ('范围：', 14, NAVY, True),
        ('单人一周集中实践（7/28–7/30），以 Obsidian 知识库 + WorkBuddy 为主要载体。', 14, DARK, False)],
       'line_spacing': 1.2}])
text(s, 0.85, 3.5, 11.6, 0.4, [{'runs': [('本周三大目标', 16, NAVY, True)]}])
goals = [
    ('建立结构化 AI 知识库', '按概念 / 工具 / 对比分类沉淀，形成可检索的体系。', BLUE),
    ('掌握 WorkBuddy 并跑通自动化', '理解核心模块，落地一个端到端自动化工作流。', ACCENT),
    ('沉淀可复用方法论', '形成可向团队复制的知识管理范式。', BLUE),
]
gx = [0.6, 4.55, 8.5]; gw = 3.6
for x, (t, d, ac) in zip(gx, goals):
    rect(s, x, 4.0, gw, 2.4, fill=BG, line=LIGHT, lw=1)
    rect(s, x, 4.0, gw, 0.10, fill=ac)
    oval(s, x + 0.25, 4.35, 0.55, 0.55, fill=ac)
    text(s, x + 0.25, 4.35, 0.55, 0.55, [{'runs': [('✓', 18, 'FFFFFF', True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.95, 4.35, gw - 1.1, 0.6, [{'runs': [(t, 15, NAVY, True)]}, {'anchor': 'mid'}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.3, 5.15, gw - 0.6, 1.1, [{'runs': [(d, 12, GRAY, False)], 'line_spacing': 1.15}])
footer(s, 4)

# =====================================================================
# Slide 5 — 周进展时间线
# =====================================================================
s = slide()
title_bar(s, '周进展时间线', 'WEEKLY TIMELINE')
days = [
    ('Day 1', '07/28', 'AI 基础概念 + 知识库搭建', '学习 LLM / Token / 上下文 / 幻觉；用 Git + Obsidian 建立 AI 概念目录结构。', 'Git 代理配置报错 → 已解决', BLUE),
    ('Day 2', '07/29', 'WorkBuddy 入门 + 工具实操', '熟悉 WorkBuddy 六大模块，绘制功能脑图；与 Claude 协作实现小恐龙跑酷游戏。', 'Claude Code 下载失败 → 关闭旧进程解决', ACCENT),
    ('Day 3', '07/30', '工作流自动化 + 方法论', '搭建 WorkBuddy 自动化工作流，每月生成知识库 HTML 可视化报告；引入 LLM Wiki 方法论。', 'Agent 与 Workflow 边界模糊 → 已厘清', BLUE),
]
# 时间轴主线
rect(s, 1.0, 3.55, 11.3, 0.04, fill=LIGHT)
for i, (d, dt, t, desc, prob, ac) in enumerate(days):
    cx = 2.3 + i * 4.4
    oval(s, cx - 0.18, 3.37, 0.4, 0.4, fill=ac)
    text(s, cx - 1.4, 2.75, 2.8, 0.5, [{'runs': [(f'{d}  {dt}', 14, ac, True)], 'align': PP_ALIGN.CENTER}])
    # 卡片
    rect(s, cx - 1.7, 4.0, 3.4, 2.6, fill=BG, line=LIGHT, lw=1)
    rect(s, cx - 1.7, 4.0, 3.4, 0.10, fill=ac)
    text(s, cx - 1.5, 4.2, 3.0, 0.5, [{'runs': [(t, 14, NAVY, True)], 'line_spacing': 1.05}])
    text(s, cx - 1.5, 4.75, 3.0, 1.0, [{'runs': [(desc, 11, GRAY, False)], 'line_spacing': 1.12}])
    text(s, cx - 1.5, 5.85, 3.0, 0.6, [{'runs': [('⚠ ' + prob, 10, ACCENT, True)], 'line_spacing': 1.05}])
footer(s, 5)

# =====================================================================
# Slide 6 — 成果一：知识库体系（含柱状图）
# =====================================================================
s = slide()
title_bar(s, '成果一：AI 知识库体系成型', 'DELIVERABLE 1 · KNOWLEDGE BASE')
cd = CategoryChartData()
cd.categories = ['AI概念', 'AI工具', '工具对比', '日报实践', 'WB实践', '资料', '索引']
cd.add_series('页数', (10, 5, 1, 5, 2, 1, 1))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.9), Inches(7.0), Inches(4.6), cd)
chart = gf.chart
chart.has_legend = False
chart.has_title = True
ct = chart.chart_title.text_frame
ct.text = '知识库各模块页数分布（共 25 页）'
ct.paragraphs[0].runs[0].font.size = Pt(13)
ct.paragraphs[0].runs[0].font.bold = True
ct.paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(NAVY)
plot = chart.plots[0]
plot.has_data_labels = True
plot.gap_width = 55
plot.data_labels.font.size = Pt(11)
plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = RGBColor.from_string(NAVY)
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = RGBColor.from_string(BLUE)
chart.category_axis.tick_labels.font.size = Pt(10)
chart.value_axis.tick_labels.font.size = Pt(9)
chart.value_axis.has_major_gridlines = True
chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor.from_string(LIGHT)
chart.value_axis.visible = False
# 右侧说明
rect(s, 8.0, 1.9, 4.73, 4.6, fill=SOFT)
rect(s, 8.0, 1.9, 0.12, 4.6, fill=ACCENT)
text(s, 8.3, 2.1, 4.3, 0.5, [{'runs': [('结构要点', 16, NAVY, True)]}])
text(s, 8.3, 2.75, 4.3, 3.5,
     [{'runs': [('• 六大分类清晰：概念 / 工具 / 对比 / 日报 / 实践 / 资料', 12, DARK, False)], 'space_after': 10, 'line_spacing': 1.1},
      {'runs': [('• 概念页 10 篇，构成知识体系主干', 12, DARK, False)], 'space_after': 10, 'line_spacing': 1.1},
      {'runs': [('• 核心枢纽：RAG（6 入链）、LLM（5 入链）', 12, DARK, False)], 'space_after': 10, 'line_spacing': 1.1},
      {'runs': [('• 索引页随每次更新自动维护，导航零成本', 12, DARK, False)], 'space_after': 10, 'line_spacing': 1.1},
      {'runs': [('• 全库为 Git 管理的 Markdown，天然带版本历史', 12, DARK, False)], 'line_spacing': 1.1}])
footer(s, 6)

# =====================================================================
# Slide 7 — 成果二：WorkBuddy 自动化工作流
# =====================================================================
s = slide()
title_bar(s, '成果二：WorkBuddy 自动化工作流', 'DELIVERABLE 2 · AUTOMATION')
steps = [
    ('原始 MD 文件', 'Obsidian 知识库\n全部 .md 文档', BLUE),
    ('扫描与统计', '目录文件数 / 总字数\nWiki 链接 / 标签分布', ACCENT),
    ('HTML 可视化报告', '图谱概览 / 笔记热度\n双链统计等关键指标', BLUE),
    ('每月 1 号自动生成', '定时调度\n零人工复盘', ACCENT),
]
bx = [0.6, 3.55, 6.5, 9.45]; bw = 2.7
for i, (t, d, ac) in enumerate(steps):
    x = bx[i]
    rect(s, x, 2.5, bw, 1.9, fill=BG, line=LIGHT, lw=1)
    rect(s, x, 2.5, bw, 0.10, fill=ac)
    text(s, x + 0.15, 2.72, bw - 0.3, 0.5, [{'runs': [(f'步骤 {i+1}', 11, ac, True)]}])
    text(s, x + 0.15, 3.12, bw - 0.3, 0.5, [{'runs': [(t, 15, NAVY, True)]}])
    text(s, x + 0.15, 3.62, bw - 0.3, 0.75, [{'runs': [(d, 11, GRAY, False)], 'line_spacing': 1.1}])
    if i < 3:
        text(s, x + bw - 0.02, 2.5, 0.45, 1.9, [{'runs': [('→', 26, NAVY, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
# 价值条
rect(s, 0.6, 4.85, 12.13, 1.5, fill=NAVY)
text(s, 0.9, 5.0, 11.5, 0.45, [{'runs': [('业务价值', 15, ACCENT, True)]}])
text(s, 0.9, 5.5, 11.6, 0.8,
     [{'runs': [
        ('• 知识复盘从手动整理升级为自动化流水线，显著降低成本；', 13, 'FFFFFF', False)], 'space_after': 4, 'line_spacing': 1.1},
      {'runs': [('• 积累"图谱 + 热度 + 双链"量化指标，为数据驱动的知识管理打底。', 13, 'FFFFFF', False)], 'line_spacing': 1.1}])
footer(s, 7)

# =====================================================================
# Slide 8 — 成果三：AI 工具实操与对比
# =====================================================================
s = slide()
title_bar(s, '成果三：AI 工具实操与对比', 'DELIVERABLE 3 · TOOLS')
# 左：实操
rect(s, 0.6, 1.85, 5.9, 4.65, fill=SOFT)
rect(s, 0.6, 1.85, 0.12, 4.65, fill=BLUE)
text(s, 0.95, 2.05, 5.4, 0.45, [{'runs': [('编码智能体实测', 16, NAVY, True)]}])
text(s, 0.95, 2.6, 5.4, 3.7,
     [{'runs': [('• Claude Code：安装与工程化协作实测', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.1},
      {'runs': [('• Codex（OpenAI CLI Agent）概览与对比', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.1},
      {'runs': [('• 完成 AI 助手横向对比：Claude / ChatGPT / Gemini / DeepSeek', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.1},
      {'runs': [('• 产出可运行 Demo：小恐龙跑酷游戏（三档速度）', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.1},
      {'runs': [('• 验证"人提需求 + Agent 协作开发"的可用性', 12, DARK, False)], 'line_spacing': 1.1}])
# 右：对比结论
rect(s, 6.7, 1.85, 6.03, 4.65, fill=BG, line=LIGHT, lw=1)
rect(s, 6.7, 1.85, 6.03, 0.10, fill=ACCENT)
text(s, 7.0, 2.05, 5.5, 0.45, [{'runs': [('关键结论', 16, NAVY, True)]}])
text(s, 7.0, 2.6, 5.5, 3.7,
     [{'runs': [('• 各助手定位不同：Claude 工程协作强，Gemini 长上下文，DeepSeek 性价比高', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• 编码 Agent 适合"脚手架 + 迭代"，关键在需求拆解与验收', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• WorkBuddy 优势在"多模块 + 连接器 + 工作流"一体化', 12, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• 选择建议：按任务类型组合使用，而非单一依赖', 12, DARK, False)], 'line_spacing': 1.12}])
footer(s, 8)

# =====================================================================
# Slide 9 — 关键数据看板
# =====================================================================
s = slide()
title_bar(s, '关键数据看板', 'KEY METRICS')
metrics = [
    ('25', '累计笔记页数', '结构化沉淀', BLUE),
    ('10', 'AI 概念覆盖', '体系主干', ACCENT),
    ('5+', '工具实操类', '含编码 Agent', BLUE),
    ('1', '自动化工作流', '月级调度', ACCENT),
    ('6', '核心枢纽入链', 'RAG 页面', BLUE),
    ('3', '实践日报', '每日复盘', ACCENT),
    ('3', '问题解决', '全部闭环', BLUE),
    ('1', '沉淀方法论', 'LLM Wiki', ACCENT),
]
xs = [0.6, 3.65, 6.7, 9.75]; cw = 2.85
ys = [1.85, 4.05]; ch = 1.95
for i, (b, l, sub, ac) in enumerate(metrics):
    x = xs[i % 4]; y = ys[i // 4]
    card(s, x, y, cw, ch, b, l, sub, ac, big_sz=40)
footer(s, 9)

# =====================================================================
# Slide 10 — 方法论：LLM Wiki
# =====================================================================
s = slide()
title_bar(s, '方法论沉淀：LLM Wiki 知识管理法', 'METHODOLOGY')
layers = [
    ('RAW SOURCES', '原始素材层', '文章 / 论文 / 数据，只读不可改，作为事实源。', BLUE),
    ('THE WIKI', '知识库层', 'LLM 生成的 Markdown 体系：摘要 / 概念 / 对比 / 索引。', ACCENT),
    ('THE SCHEMA', '规范层', '约定结构、命名与工作流（如 AGENTS.md），让 LLM 成为守纪律的维护者。', BLUE),
]
lx = [0.6, 4.85, 9.1]; lw = 3.5
for x, (en, cn, d, ac) in zip(lx, layers):
    rect(s, x, 1.95, lw, 2.5, fill=BG, line=LIGHT, lw=1)
    rect(s, x, 1.95, lw, 0.10, fill=ac)
    text(s, x + 0.2, 2.15, lw - 0.4, 0.35, [{'runs': [(en, 12, ac, True)]}])
    text(s, x + 0.2, 2.5, lw - 0.4, 0.45, [{'runs': [(cn, 18, NAVY, True)]}])
    text(s, x + 0.2, 3.05, lw - 0.4, 1.3, [{'runs': [(d, 12, GRAY, False)], 'line_spacing': 1.15}])
    if x != lx[-1]:
        text(s, x + lw - 0.02, 1.95, 0.45, 2.5, [{'runs': [('→', 26, NAVY, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
rect(s, 0.6, 4.75, 12.13, 1.6, fill=SOFT)
rect(s, 0.6, 4.75, 0.12, 1.6, fill=ACCENT)
text(s, 0.95, 4.95, 11.5, 0.45, [{'runs': [('核心理念：知识可"复利"', 15, NAVY, True)]}])
text(s, 0.95, 5.45, 11.6, 0.85,
     [{'runs': [('传统 RAG 每次提问都从头检索、无积累；LLM Wiki 由 LLM 增量构建并持续维护持久知识库——交叉引用已就位、矛盾已标注、综合已成型，越用越厚。适用于个人、研究、团队内部 wiki 与竞品分析。', 12, DARK, False)], 'line_spacing': 1.15}])
footer(s, 10)

# =====================================================================
# Slide 11 — 问题与解决
# =====================================================================
s = slide()
title_bar(s, '问题与解决', 'ISSUES & RESOLUTION')
items = [
    ('Git 配置多次报错，Obsidian 无法实时上传', '排查代理配置，修正后上传链路恢复。', BLUE),
    ('Claude Code 下载失败', '原因为旧进程未关闭，结束进程后下载成功。', ACCENT),
    ('Agent 与 Workflow 边界模糊', '查阅官方文档 + 实际调试，明确 Workflow 为执行载体、Agent 为其中智能组件，据此重构。', BLUE),
    ('知识库结构初期零散', '引入索引页与双链规范，按概念轴心组织语义网络。', ACCENT),
]
y = 1.85
for i, (p, sol, ac) in enumerate(items):
    ry = y + i * 1.18
    rect(s, 0.6, ry, 12.13, 1.0, fill=SOFT)
    rect(s, 0.6, ry, 0.12, 1.0, fill=ac)
    text(s, 0.95, ry + 0.12, 7.3, 0.8,
         [{'runs': [('问题  ', 11, ac, True), (p, 13, NAVY, True)], 'line_spacing': 1.05}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.4, ry + 0.12, 4.2, 0.8,
         [{'runs': [('解决  ', 11, '2E8B57', True), (sol, 12, DARK, False)], 'line_spacing': 1.05}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11)

# =====================================================================
# Slide 12 — 价值与影响
# =====================================================================
s = slide()
title_bar(s, '价值与影响', 'VALUE & IMPACT')
# 左：个人
rect(s, 0.6, 1.85, 5.9, 4.65, fill=BG, line=LIGHT, lw=1)
rect(s, 0.6, 1.85, 5.9, 0.10, fill=BLUE)
text(s, 0.95, 2.1, 5.4, 0.45, [{'runs': [('对个人', 17, NAVY, True)]}])
text(s, 0.95, 2.7, 5.4, 3.6,
     [{'runs': [('• 建立 AI 与编码 Agent 的系统认知', 13, DARK, False)], 'space_after': 10, 'line_spacing': 1.12},
      {'runs': [('• 掌握 WorkBuddy 核心模块与自动化能力', 13, DARK, False)], 'space_after': 10, 'line_spacing': 1.12},
      {'runs': [('• 形成可检索、可演进的个人知识资产', 13, DARK, False)], 'space_after': 10, 'line_spacing': 1.12},
      {'runs': [('• 复盘效率从"手动整理"转向"自动生成"', 13, DARK, False)], 'line_spacing': 1.12}])
# 右：团队（建议）
rect(s, 6.7, 1.85, 6.03, 4.65, fill=SOFT)
rect(s, 6.7, 1.85, 6.03, 0.10, fill=ACCENT)
text(s, 7.05, 2.1, 5.5, 0.45, [{'runs': [('对团队（建议推广）', 17, NAVY, True)]}])
text(s, 7.05, 2.7, 5.5, 3.0,
     [{'runs': [('• 可直接复用的知识库搭建模板', 13, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• 自动化复盘工作流可迁移到团队 wiki', 13, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• LLM Wiki 方法论支撑持续知识沉淀', 13, DARK, False)], 'space_after': 9, 'line_spacing': 1.12},
      {'runs': [('• 降低新人上手 AI 工具的学习曲线', 13, DARK, False)], 'line_spacing': 1.12}])
text(s, 6.7, 5.95, 6.03, 0.5,
     [{'runs': [('* 量化收益待团队试点后测算（方向性预估）', 10, GRAY, False)], 'align': PP_ALIGN.CENTER}])
footer(s, 12)

# =====================================================================
# Slide 13 — 下一步计划
# =====================================================================
s = slide()
title_bar(s, '下一步计划', 'NEXT STEPS')
phases = [
    ('短期 · 本周', '知识库结构化升级', ['概念驱动 + 双链强化', '消除知识断点', '补全场景与案例'], BLUE),
    ('中期 · 1 个月', '团队试点自动化', ['选 1 个团队知识库场景', '跑通月度可视化报告', '收集反馈迭代'], ACCENT),
    ('长期 · 季度', '沉淀使用规范 SOP', ['形成 AI 工具使用规范', '建立复盘与评审机制', '规模化复制'], BLUE),
]
px = [0.6, 4.85, 9.1]; pw = 3.5
for x, (when, t, pts, ac) in zip(px, phases):
    rect(s, x, 1.9, pw, 4.6, fill=BG, line=LIGHT, lw=1)
    rect(s, x, 1.9, pw, 0.55, fill=ac)
    text(s, x, 1.9, pw, 0.55, [{'runs': [(when, 14, 'FFFFFF', True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.25, 2.65, pw - 0.5, 0.5, [{'runs': [(t, 16, NAVY, True)]}])
    paras = [{'runs': [('• ' + p, 13, DARK, False)], 'space_after': 12, 'line_spacing': 1.12} for p in pts]
    text(s, x + 0.25, 3.35, pw - 0.5, 2.9, paras)
footer(s, 13)

# =====================================================================
# Slide 14 — 风险与所需支持
# =====================================================================
s = slide()
title_bar(s, '风险与所需支持', 'RISKS & SUPPORT')
# 左 风险
rect(s, 0.6, 1.85, 5.9, 4.65, fill=BG, line=LIGHT, lw=1)
rect(s, 0.6, 1.85, 0.12, 4.65, fill="C0392B")
text(s, 0.95, 2.05, 5.4, 0.45, [{'runs': [('主要风险', 16, 'C0392B', True)]}])
text(s, 0.95, 2.65, 5.4, 3.6,
     [{'runs': [('• 当前为单人实践，尚未规模化验证', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 数据以基础沉淀为主，ROI 待量化', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 自动化工作流依赖特定环境与连接器', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 方法论推广需配套使用习惯培养', 13, DARK, False)], 'line_spacing': 1.12}])
# 右 支持
rect(s, 6.7, 1.85, 6.03, 4.65, fill=SOFT)
rect(s, 6.7, 1.85, 0.12, 4.65, fill=BLUE)
text(s, 7.05, 2.05, 5.5, 0.45, [{'runs': [('所需支持', 16, NAVY, True)]}])
text(s, 7.05, 2.65, 5.5, 3.6,
     [{'runs': [('• 1 个团队知识库场景的试点授权', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 试点所需的连接器和环境资源', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 方法论评审与推广节奏确认', 13, DARK, False)], 'space_after': 11, 'line_spacing': 1.12},
      {'runs': [('• 月度复盘报告的受众与反馈机制', 13, DARK, False)], 'line_spacing': 1.12}])
footer(s, 14)

# =====================================================================
# Slide 15 — 总结（一句话结论）
# =====================================================================
s = slide()
rect(s, 0, 0, W, H, fill=NAVY)
rect(s, 0, 0, 0.28, H, fill=ACCENT)
oval(s, 10.6, 4.6, 3.4, 3.4, fill="27406A")
text(s, 1.0, 2.0, 11.4, 0.5, [{'runs': [('总结 · CONCLUSION', 14, ACCENT, True)]}])
text(s, 1.0, 2.7, 11.4, 2.2,
     [{'runs': [('一周搭建起"AI 知识库 + 自动化复盘"的最小可行闭环，', 30, 'FFFFFF', True)], 'line_spacing': 1.25, 'space_after': 6},
      {'runs': [('已具备向团队复制推广的条件。', 30, 'FFFFFF', True)], 'line_spacing': 1.25}],
     anchor=MSO_ANCHOR.MIDDLE)
rect(s, 1.0, 5.2, 3.0, 0.06, fill=ACCENT)
text(s, 1.0, 5.5, 11.4, 0.6, [{'runs': [('建议：批准一个团队场景试点，验证规模化价值。', 15, 'C9D6E5', False)]}])

# =====================================================================
# Slide 16 — 结束页
# =====================================================================
s = slide()
rect(s, 0, 0, W, H, fill=BG)
oval(s, -1.2, -1.2, 3.4, 3.4, fill=SOFT)
oval(s, 11.4, 5.2, 3.0, 3.0, fill="EAF1F8")
text(s, 0, 2.7, W, 1.0, [{'runs': [('谢谢观看 · 欢迎讨论', 40, NAVY, True)], 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
rect(s, W/2 - 1.5, 3.85, 3.0, 0.06, fill=ACCENT)
text(s, 0, 4.2, W, 0.6,
     [{'runs': [('第一周项目汇报  ·  AI 工具上手与 WorkBuddy 自动化实践', 14, GRAY, False)], 'align': PP_ALIGN.CENTER}])
text(s, 0, 4.85, W, 0.5,
     [{'runs': [('汇报日期：2026/07/31', 13, GRAY, False)], 'align': PP_ALIGN.CENTER}])

# ---------- 输出 ----------
out = r'D:\obsidian1\第一周\项目汇报PPT\第一周项目汇报.pptx'
prs.save(out)
print('SAVED', out, 'slides=', len(prs.slides._sldIdLst))
